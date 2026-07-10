import replicate
import traceback
import os
import json
import uuid
import base64
import io
import csv
import xml.etree.ElementTree as ET
from datetime import datetime
import streamlit as st
from sentence_transformers import SentenceTransformer
import chromadb
from groq import Groq
import fitz
from pptx import Presentation
from docx import Document
import openpyxl

st.set_page_config(page_title="Aquanis", page_icon="💧", layout="wide")

# ---------- Design tokens ----------
BG = "#03111e"
FG = "#e7f0f5"
CARD = "#0a1c2c"
PRIMARY = "#3dbfe2"
PRIMARY_FG = "#010e1d"
SECONDARY = "#112b40"
MUTED_FG = "#879ca8"
ACCENT = "#10364e"
BORDER = "rgba(119, 184, 215, 0.15)"
SIDEBAR = "#051729"
SIDEBAR_ACCENT = "#0f283d"

TRANSLATIONS = {
    "en": {
        "welcome_title": "Welcome to Aquanis",
        "welcome_caption": "Sign in to save your chat history, or continue without an account.",
        "sign_in": "Sign in with Google",
        "continue_guest": "Continue without signing in",
        "guest_warning": "Guest chats are not saved and will be lost if you refresh the page.",
        "new_chat": "+ New chat",
        "chat_name_label": "Chat name",
        "chat_name_placeholder": "e.g. Reynolds number",
        "create": "Create",
        "cancel": "Cancel",
        "recent_chats": "Recent chats",
        "user_label": "User",
        "guest_label": "Guest",
        "log_out": "Log out",
        "app_title": "New chat",
        "welcome_sub": "Your AI companion for fluid mechanics and hydraulic engineering. Ask about pipe flow, pumps, open channels, hydrology, and more.",
        "chat_input_placeholder": "Ask a question about hydraulics, or attach a file...",
        "thinking": "Aquanis is thinking...",
        "sources_label": "Sources",
        "language_label": "Interface language",
        "footer_note": "Aquanis can make mistakes. Verify critical hydraulic calculations.",
        "suggestions": [
            "Explain the Bernoulli equation with a practical example",
            "How do I size a centrifugal pump for a pipeline?",
            "Derive the Darcy-Weisbach head loss formula",
            "What causes water hammer and how do I prevent it?",
        ],
    },
    "fr": {
        "welcome_title": "Bienvenue sur Aquanis",
        "welcome_caption": "Connectez-vous pour sauvegarder vos discussions, ou continuez sans compte.",
        "sign_in": "Se connecter avec Google",
        "continue_guest": "Continuer sans se connecter",
        "guest_warning": "Les discussions invité ne sont pas sauvegardées et seront perdues si vous actualisez la page.",
        "new_chat": "+ Nouvelle discussion",
        "chat_name_label": "Nom de la discussion",
        "chat_name_placeholder": "ex: Nombre de Reynolds",
        "create": "Créer",
        "cancel": "Annuler",
        "recent_chats": "Discussions récentes",
        "user_label": "Utilisateur",
        "guest_label": "Invité",
        "log_out": "Se déconnecter",
        "app_title": "Nouvelle discussion",
        "welcome_sub": "Votre assistant IA pour la mécanique des fluides et l'hydraulique. Posez vos questions sur les écoulements, les pompes, les canaux ouverts, l'hydrologie, et plus encore.",
        "chat_input_placeholder": "Posez une question, ou joignez un fichier...",
        "thinking": "Aquanis réfléchit...",
        "sources_label": "Sources",
        "language_label": "Langue de l'interface",
        "footer_note": "Aquanis peut faire des erreurs. Vérifiez les calculs hydrauliques critiques.",
        "suggestions": [
            "Expliquer l'équation de Bernoulli avec un exemple pratique",
            "Comment dimensionner une pompe centrifuge pour une conduite ?",
            "Démontrer la formule de perte de charge de Darcy-Weisbach",
            "Quelles sont les causes du coup de bélier et comment l'éviter ?",
        ],
    },
    "ar": {
        "welcome_title": "مرحبا بك في Aquanis",
        "welcome_caption": "سجل الدخول لحفظ محادثاتك، أو تابع بدون حساب.",
        "sign_in": "تسجيل الدخول بجوجل",
        "continue_guest": "المتابعة بدون تسجيل الدخول",
        "guest_warning": "لا يتم حفظ محادثات الضيف وستفقد عند تحديث الصفحة.",
        "new_chat": "+ محادثة جديدة",
        "chat_name_label": "اسم المحادثة",
        "chat_name_placeholder": "مثال: عدد رينولدز",
        "create": "إنشاء",
        "cancel": "إلغاء",
        "recent_chats": "المحادثات الأخيرة",
        "user_label": "المستخدم",
        "guest_label": "ضيف",
        "log_out": "تسجيل الخروج",
        "app_title": "محادثة جديدة",
        "welcome_sub": "مساعدك الذكي في ميكانيكا الموائع والهندسة الهيدروليكية. اسأل عن جريان الأنابيب، المضخات، القنوات المفتوحة، الهيدرولوجيا، والمزيد.",
        "chat_input_placeholder": "اطرح سؤالا، أو أرفق ملفا...",
        "thinking": "Aquanis يفكر...",
        "sources_label": "المصادر",
        "language_label": "لغة الواجهة",
        "footer_note": "قد يخطئ Aquanis. تحقق من الحسابات الهيدروليكية الهامة.",
        "suggestions": [
            "اشرح معادلة برنولي بمثال عملي",
            "كيف أحدد حجم مضخة طاردة مركزية لخط أنابيب؟",
            "اشتق معادلة فقدان الضغط دارسي-فايسباخ",
            "ما أسباب المطرقة المائية وكيف أمنعها؟",
        ],
    },
}

if "ui_lang" not in st.session_state:
    st.session_state.ui_lang = "en"

t = TRANSLATIONS[st.session_state.ui_lang]
is_rtl = st.session_state.ui_lang == "ar"

if "guest_mode" not in st.session_state:
    st.session_state.guest_mode = False
if "guest_id" not in st.session_state:
    st.session_state.guest_id = "guest_" + str(uuid.uuid4())

# ---------- Global theme CSS ----------
st.markdown(f"""
<style>
.stApp {{
    background-color: {BG};
    color: {FG};
    {"direction: rtl;" if is_rtl else ""}
}}
[data-testid="stSidebar"] {{
    background-color: {SIDEBAR};
    border-right: 1px solid {BORDER};
}}
[data-testid="stSidebar"] * {{
    color: {FG} !important;
}}
.stButton button {{
    background-color: transparent;
    border: 1px solid transparent;
    text-align: left;
    color: {FG};
    border-radius: 10px;
}}
.stButton button:hover {{
    background-color: {SIDEBAR_ACCENT};
    color: {PRIMARY};
    border: 1px solid {BORDER};
}}
[data-testid="stChatInput"] {{
    background-color: {CARD};
    border: 1px solid {BORDER};
    border-radius: 16px;
}}
.aquanis-logo {{
    display: inline-flex;
    align-items: center;
    justify-content: center;
    width: 40px;
    height: 40px;
    border-radius: 50%;
    background-color: rgba(61, 191, 226, 0.15);
    font-size: 20px;
    margin-right: 8px;
}}
.aquanis-header {{
    display: flex;
    align-items: center;
    border-bottom: 1px solid {BORDER};
    padding-bottom: 14px;
    margin-bottom: 8px;
}}
.aquanis-user-bubble {{
    display: flex;
    justify-content: flex-end;
    margin: 10px 0;
}}
.aquanis-user-bubble-inner {{
    max-width: 75%;
    background-color: {PRIMARY};
    color: {PRIMARY_FG};
    padding: 12px 16px;
    border-radius: 18px 18px 4px 18px;
    font-size: 14px;
    line-height: 1.5;
}}
.aquanis-assistant-bubble {{
    display: flex;
    align-items: flex-start;
    gap: 10px;
    margin: 10px 0;
}}
.aquanis-assistant-bubble-inner {{
    max-width: 75%;
    background-color: {CARD};
    border: 1px solid {BORDER};
    color: {FG};
    padding: 12px 16px;
    border-radius: 18px 18px 18px 4px;
    font-size: 14px;
    line-height: 1.5;
    white-space: pre-line;
}}
.aquanis-suggestion {{
    border: 1px solid {BORDER};
    background-color: {CARD};
    border-radius: 14px;
    padding: 14px;
    font-size: 13px;
    color: {FG};
}}
.aquanis-footer-note {{
    text-align: center;
    font-size: 11px;
    color: {MUTED_FG};
    margin-top: 8px;
}}
</style>
""", unsafe_allow_html=True)


# ---------- File extraction functions ----------
def extract_text_from_pdf(file_bytes):
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    return "\n".join(page.get_text() for page in doc)


def extract_text_from_pptx(file_bytes):
    prs = Presentation(io.BytesIO(file_bytes))
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                parts.append(shape.text)
    return "\n".join(parts)


def extract_text_from_docx(file_bytes):
    doc = Document(io.BytesIO(file_bytes))
    return "\n".join(p.text for p in doc.paragraphs)


def extract_text_from_xlsx(file_bytes):
    wb = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True)
    parts = []
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            row_text = " | ".join(str(cell) for cell in row if cell is not None)
            if row_text:
                parts.append(row_text)
    return "\n".join(parts)
# ---------- Image generation function ----------
def generate_diagram(prompt):
    """Generate a diagram using Stable Diffusion via Replicate"""
    try:
        output = replicate_client.run(
            "stability-ai/sdxl:39ed52f2a60c3b36b4e8c8cb03d33ec2e7925ea2b3b9a44cc27a992cb5d52e27",
            input={
                "prompt": prompt + " technical diagram, engineering drawing, clear labels, professional",
                "num_inference_steps": 25,
                "guidance_scale": 7.5,
            }
        )
        if output and len(output) > 0:
            return output[0]
    except Exception as e:
        print(f"Image generation error: {e}")
    return None

def extract_text_from_csv(file_bytes):
    text_lines = []
    try:
        for row in csv.reader(file_bytes.decode('utf-8').splitlines()):
            text_lines.append(" | ".join(row))
    except:
        text_lines.append(file_bytes.decode('utf-8'))
    return "\n".join(text_lines)


def extract_text_from_txt(file_bytes):
    return file_bytes.decode('utf-8', errors='ignore')


def extract_text_from_json(file_bytes):
    try:
        data = json.loads(file_bytes.decode('utf-8'))
        return json.dumps(data, indent=2)
    except:
        return file_bytes.decode('utf-8', errors='ignore')


def extract_text_from_xml(file_bytes):
    try:
        root = ET.fromstring(file_bytes)
        return ET.tostring(root, encoding='unicode')
    except:
        return file_bytes.decode('utf-8', errors='ignore')


def extract_text_from_rtf(file_bytes):
    text = file_bytes.decode('utf-8', errors='ignore')
    import re
    text = re.sub(r'\\[a-z]+\d*\s?', '', text)
    text = re.sub(r'[{}]', '', text)
    return text


try:
    is_logged_in = st.user.is_logged_in

    if not is_logged_in and not st.session_state.guest_mode:
        st.markdown("<div class='aquanis-logo'>💧</div>", unsafe_allow_html=True)
        st.markdown("### " + t["welcome_title"])
        st.caption(t["welcome_caption"])
        col1, col2 = st.columns(2)
        with col1:
            if st.button(t["sign_in"], use_container_width=True):
                st.login()
        with col2:
            if st.button(t["continue_guest"], use_container_width=True):
                st.session_state.guest_mode = True
                st.rerun()
        st.caption(t["guest_warning"])
        st.stop()

    if is_logged_in:
        user_identity = st.user.email
        display_name = st.user.name
    else:
        user_identity = st.session_state.guest_id
        display_name = t["guest_label"]

    groq_client = Groq(api_key=st.secrets["GROQ_API_KEY"])
    CHROMA_PATH = "chroma_db"
    CHATS_FILE = "chats.json"
    replicate_client = replicate.Client(api_token=st.secrets.get("REPLICATE_API_KEY"))

    def load_all_chats():
        if os.path.exists(CHATS_FILE):
            with open(CHATS_FILE, "r", encoding="utf-8") as f:
                return json.load(f)
        return {}

    def save_all_chats(all_chats):
        with open(CHATS_FILE, "w", encoding="utf-8") as f:
            json.dump(all_chats, f, ensure_ascii=False, indent=2)

    def load_chats(user_key):
        all_chats = load_all_chats()
        return all_chats.get(user_key, {})

    def save_chats(user_key, chats):
        all_chats = load_all_chats()
        all_chats[user_key] = chats
        save_all_chats(all_chats)

    if "chats" not in st.session_state:
        st.session_state.chats = load_chats(user_identity) if is_logged_in else {}

    if "current_chat_id" not in st.session_state:
        st.session_state.current_chat_id = (
            list(st.session_state.chats.keys())[0] if st.session_state.chats else None
        )

    if "creating_new_chat" not in st.session_state:
        st.session_state.creating_new_chat = False

    @st.cache_resource
    def load_resources():
        model = SentenceTransformer("all-MiniLM-L6-v2")
        client = chromadb.PersistentClient(path=CHROMA_PATH)
        collection = client.get_collection("aquanis_docs")
        return model, collection

    model, collection = load_resources()

    # ---------- Sidebar ----------
    with st.sidebar:
        st.markdown(
            f"<div style='display:flex; align-items:center; gap:10px; padding:4px 0 12px;'>"
            f"<span class='aquanis-logo'>💧</span>"
            f"<div><div style='font-weight:600; font-size:17px;'>Aquanis</div>"
            f"<div style='font-size:12px; color:{MUTED_FG};'>🌊 Hydraulics AI</div></div></div>",
            unsafe_allow_html=True
        )

        lang_options = {"English": "en", "Français": "fr", "العربية": "ar"}
        lang_names = list(lang_options.keys())
        current_lang_name = [k for k, v in lang_options.items() if v == st.session_state.ui_lang][0]
        selected_lang_name = st.selectbox(t["language_label"], lang_names, index=lang_names.index(current_lang_name))
        selected_lang_code = lang_options[selected_lang_name]
        if selected_lang_code != st.session_state.ui_lang:
            st.session_state.ui_lang = selected_lang_code
            st.rerun()

        if st.button(t["new_chat"], use_container_width=True, type="primary"):
            st.session_state.creating_new_chat = True
            st.rerun()

        if st.session_state.creating_new_chat:
            new_name = st.text_input(t["chat_name_label"], placeholder=t["chat_name_placeholder"], key="new_chat_name")
            col_a, col_b = st.columns(2)
            with col_a:
                if st.button(t["create"], use_container_width=True):
                    new_id = str(uuid.uuid4())
                    title = new_name.strip() if new_name.strip() else t["new_chat"]
                    st.session_state.chats[new_id] = {
                        "title": title, "messages": [], "created": datetime.now().isoformat()
                    }
                    st.session_state.current_chat_id = new_id
                    st.session_state.creating_new_chat = False
                    if is_logged_in:
                        save_chats(user_identity, st.session_state.chats)
                    st.rerun()
            with col_b:
                if st.button(t["cancel"], use_container_width=True):
                    st.session_state.creating_new_chat = False
                    st.rerun()

        st.markdown("---")
        st.markdown(
            f"<p style='font-size:11px; letter-spacing:0.5px; color:{MUTED_FG}; text-transform:uppercase; margin-bottom:8px;'>{t['recent_chats']}</p>",
            unsafe_allow_html=True
        )

        for chat_id, chat in sorted(st.session_state.chats.items(), key=lambda x: x[1]["created"], reverse=True):
            label = "💬 " + (chat["title"] if chat["title"] else t["new_chat"])
            col1, col2 = st.columns([5, 1])
            with col1:
                if st.button(label, key=f"chat_{chat_id}", use_container_width=True):
                    st.session_state.current_chat_id = chat_id
                    st.rerun()
            with col2:
                if st.button("\U0001F5D1", key=f"del_{chat_id}"):
                    del st.session_state.chats[chat_id]
                    if st.session_state.current_chat_id == chat_id:
                        remaining = list(st.session_state.chats.keys())
                        st.session_state.current_chat_id = remaining[0] if remaining else None
                    if is_logged_in:
                        save_chats(user_identity, st.session_state.chats)
                    st.rerun()

        st.markdown("---")
        role_label = t["guest_label"] if not is_logged_in else "Student"
        col_a, col_b = st.columns([5, 1])
        with col_a:
            st.markdown(
                f"<div style='display:flex; align-items:center; gap:10px;'>"
                f"<span style='display:flex; align-items:center; justify-content:center; width:34px; height:34px; "
                f"border-radius:50%; background-color:rgba(61,191,226,0.15); font-size:16px;'>👤</span>"
                f"<div><div style='font-size:13px; font-weight:500;'>{display_name}</div>"
                f"<div style='font-size:11px; color:{MUTED_FG};'>{role_label}</div></div></div>",
                unsafe_allow_html=True
            )
        with col_b:
            if is_logged_in:
                if st.button("↪", key="logout_btn"):
                    st.logout()
            else:
                if st.button("↪", key="logout_btn"):
                    st.session_state.guest_mode = False
                    st.session_state.chats = {}
                    st.session_state.current_chat_id = None
                    st.rerun()

    # ---------- Main area ----------
    current_id = st.session_state.current_chat_id

    if current_id is None:
        st.markdown(
            "<div class='aquanis-header'><span class='aquanis-logo'>💧</span>"
            "<span style='font-size:18px; font-weight:600;'>" + t["app_title"] + "</span></div>",
            unsafe_allow_html=True
        )
        st.markdown("<div style='text-align:center; margin-top:40px;'><span style='font-size:48px;'>💧</span></div>", unsafe_allow_html=True)
        st.markdown(f"<h3 style='text-align:center;'>{t['welcome_title']}</h3>", unsafe_allow_html=True)
        st.markdown(f"<p style='text-align:center; color:{MUTED_FG}; max-width:600px; margin:0 auto 24px;'>{t['welcome_sub']}</p>", unsafe_allow_html=True)

        cols = st.columns(2)
        clicked_suggestion = None
        for i, suggestion in enumerate(t["suggestions"]):
            with cols[i % 2]:
                if st.button(suggestion, key=f"sugg_{i}", use_container_width=True):
                    clicked_suggestion = suggestion

        if clicked_suggestion:
            new_id = str(uuid.uuid4())
            st.session_state.chats[new_id] = {
                "title": clicked_suggestion[:40], "messages": [], "created": datetime.now().isoformat()
            }
            st.session_state.current_chat_id = new_id
            if is_logged_in:
                save_chats(user_identity, st.session_state.chats)
            st.session_state["pending_question"] = clicked_suggestion
            st.rerun()

        st.stop()

    current_chat = st.session_state.chats[current_id]

    st.markdown(
        "<div class='aquanis-header'><span class='aquanis-logo'>💧</span>"
        "<span style='font-size:18px; font-weight:600;'>" + current_chat["title"] + "</span></div>",
        unsafe_allow_html=True
    )

    for msg in current_chat["messages"]:
        if msg["role"] == "user":
            st.markdown(
                f"<div class='aquanis-user-bubble'><div class='aquanis-user-bubble-inner'>{msg['content']}</div></div>",
                unsafe_allow_html=True
            )
        else:
            st.markdown(
                f"<div class='aquanis-assistant-bubble'><span class='aquanis-logo'>💧</span>"
                f"<div class='aquanis-assistant-bubble-inner'>{msg['content']}</div></div>",
                unsafe_allow_html=True
            )

    prompt = st.chat_input(
        t["chat_input_placeholder"],
        accept_file=True,
        file_type=["png", "jpg", "jpeg", "pdf", "docx", "pptx", "xlsx", "xls", "csv", "txt", "json", "xml", "rtf"]
    )

    pending = st.session_state.pop("pending_question", None)
    if pending and not prompt:
        class FakePrompt:
            text = pending
            def __getitem__(self, key):
                return []
        prompt = FakePrompt()

    if prompt:
        question = prompt.text if prompt.text else ""
        uploaded_files = prompt["files"] if prompt["files"] else []

        image_data_url = None
        extra_text_context = ""

        for f in uploaded_files:
            file_bytes = f.read()
            ext = f.name.split(".")[-1].lower()

            try:
                if ext in ["png", "jpg", "jpeg"]:
                    base64_image = base64.b64encode(file_bytes).decode("utf-8")
                    image_data_url = "data:" + f.type + ";base64," + base64_image
                elif ext == "pdf":
                    extra_text_context += "\n\n[Content from " + f.name + "]\n" + extract_text_from_pdf(file_bytes)
                elif ext == "pptx":
                    extra_text_context += "\n\n[Content from " + f.name + "]\n" + extract_text_from_pptx(file_bytes)
                elif ext == "docx":
                    extra_text_context += "\n\n[Content from " + f.name + "]\n" + extract_text_from_docx(file_bytes)
                elif ext in ["xlsx", "xls"]:
                    extra_text_context += "\n\n[Content from " + f.name + "]\n" + extract_text_from_xlsx(file_bytes)
                elif ext == "csv":
                    extra_text_context += "\n\n[Content from " + f.name + "]\n" + extract_text_from_csv(file_bytes)
                elif ext == "txt":
                    extra_text_context += "\n\n[Content from " + f.name + "]\n" + extract_text_from_txt(file_bytes)
                elif ext == "json":
                    extra_text_context += "\n\n[Content from " + f.name + "]\n" + extract_text_from_json(file_bytes)
                elif ext == "xml":
                    extra_text_context += "\n\n[Content from " + f.name + "]\n" + extract_text_from_xml(file_bytes)
                elif ext == "rtf":
                    extra_text_context += "\n\n[Content from " + f.name + "]\n" + extract_text_from_rtf(file_bytes)
            except Exception as e:
                extra_text_context += f"\n\n[Could not extract content from {f.name}: {str(e)}]"

        display_text = question if question else "(file attached)"
        current_chat["messages"].append({"role": "user", "content": display_text})
        st.markdown(
            f"<div class='aquanis-user-bubble'><div class='aquanis-user-bubble-inner'>{display_text}</div></div>",
            unsafe_allow_html=True
        )

        query_embedding = model.encode([display_text]).tolist()
        results = collection.query(query_embeddings=query_embedding, n_results=4)
        context = "\n\n".join(results["documents"][0])
        sources = list(set(r["source"] for r in results["metadatas"][0]))

        system_prompt = ("You are Aquanis, a helpful assistant for hydraulics engineers and students. "
                          "Always answer in the same language the student used in their latest question. "
                          "Use the course context below to answer questions. If an image or file is attached, "
                          "analyze it and relate it to hydraulics concepts. Always write mathematical equations and "
                          "formulas using LaTeX syntax with single $ for inline and $$ for standalone equations. "
                          "When a student asks you to draw, sketch, diagram, or visualize something related to hydraulics, "
                          "respond with: [GENERATE_IMAGE: description of what to draw]. "
                          "If the answer is not available, say so in the student's language. Use earlier conversation for follow-ups.\n\n"
                          "Course context:\n" + context)

        if extra_text_context:
            system_prompt += "\n\nAttached file content:\n" + extra_text_context

        conversation_messages = [{"role": "system", "content": system_prompt}]
        num_messages = len(current_chat["messages"])
        for i, msg in enumerate(current_chat["messages"]):
            if i == num_messages - 1 and image_data_url:
                conversation_messages.append({
                    "role": "user",
                    "content": [
                        {"type": "text", "text": msg["content"]},
                        {"type": "image_url", "image_url": {"url": image_data_url}}
                    ]
                })
            else:
                conversation_messages.append({"role": msg["role"], "content": msg["content"]})

        model_to_use = "meta-llama/llama-4-maverick-17b-128e-instruct" if image_data_url else "llama-3.3-70b-versatile"

        with st.spinner(t["thinking"]):
            response = groq_client.chat.completions.create(model=model_to_use, messages=conversation_messages)
            answer = response.choices[0].message.content

        # Check if the response asks to generate an image
        if "[GENERATE_IMAGE:" in answer:
            parts = answer.split("[GENERATE_IMAGE:")
            text_part = parts[0].strip()
            image_prompt_part = parts[1].split("]")[0].strip() if len(parts) > 1 else ""
            remaining_text = parts[1].split("]", 1)[1].strip() if len(parts) > 1 and "]" in parts[1] else ""

            # Display text
            if text_part:
                st.markdown(
                    f"<div class='aquanis-assistant-bubble'><span class='aquanis-logo'>💧</span>"
                    f"<div class='aquanis-assistant-bubble-inner'>{text_part}</div></div>",
                    unsafe_allow_html=True
                )

            # Generate and display image
            if image_prompt_part:
                with st.spinner("🎨 Generating diagram..."):
                    image_url = generate_diagram(image_prompt_part)
                    if image_url:
                        st.image(image_url, caption=image_prompt_part, use_container_width=True)

            # Display remaining text
            if remaining_text:
                final_answer = remaining_text + "\n\n" + t["sources_label"] + ": " + ", ".join(sources)
                st.markdown(
                    f"<div class='aquanis-assistant-bubble'><span class='aquanis-logo'>💧</span>"
                    f"<div class='aquanis-assistant-bubble-inner'>{final_answer}</div></div>",
                    unsafe_allow_html=True
                )
            else:
                st.markdown(
                    f"<div class='aquanis-assistant-bubble'><span class='aquanis-logo'>💧</span>"
                    f"<div class='aquanis-assistant-bubble-inner'>{t['sources_label']}: {', '.join(sources)}</div></div>",
                    unsafe_allow_html=True
                )

            answer = text_part + "\n[Diagram generated]" + remaining_text
        else:
            answer = answer + "\n\n" + t["sources_label"] + ": " + ", ".join(sources)
            st.markdown(
                f"<div class='aquanis-assistant-bubble'><span class='aquanis-logo'>💧</span>"
                f"<div class='aquanis-assistant-bubble-inner'>{answer}</div></div>",
                unsafe_allow_html=True
            )

        current_chat["messages"].append({"role": "assistant", "content": answer})
        if is_logged_in:
            save_chats(user_identity, st.session_state.chats)
        st.rerun()

    st.markdown(f"<p class='aquanis-footer-note'>{t['footer_note']}</p>", unsafe_allow_html=True)

except Exception as e:
    st.error("An error occurred while running Aquanis:")
    st.code(traceback.format_exc())