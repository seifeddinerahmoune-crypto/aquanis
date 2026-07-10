import os
import csv
import json
import xml.etree.ElementTree as ET
import re
from sentence_transformers import SentenceTransformer
import chromadb
import fitz
from pptx import Presentation
from docx import Document
import openpyxl
from PIL import Image
import pytesseract

DATA_FOLDER = "data"
CHROMA_PATH = "chroma_db"

# Configure Tesseract path (Windows users - update this if Tesseract is in a different location)
try:
    pytesseract.pytesseract.pytesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'
except:
    pass

print("Loading embedding model...")
model = SentenceTransformer("all-MiniLM-L6-v2")

client = chromadb.PersistentClient(path=CHROMA_PATH)
collection = client.get_or_create_collection("aquanis_docs")

def chunk_text(text, chunk_size=400):
    words = text.split()
    return [" ".join(words[i:i+chunk_size]) for i in range(0, len(words), chunk_size)]

def read_pdf(path):
    doc = fitz.open(path)
    return "\n".join(page.get_text() for page in doc)

def read_pptx(path):
    prs = Presentation(path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def read_docx(path):
    doc = Document(path)
    return "\n".join(p.text for p in doc.paragraphs)

def read_xlsx(path):
    wb = openpyxl.load_workbook(path, data_only=True)
    parts = []
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            row_text = " | ".join(str(cell) for cell in row if cell is not None)
            if row_text:
                parts.append(row_text)
    return "\n".join(parts)

def read_csv(path):
    text_lines = []
    try:
        with open(path, 'r', encoding='utf-8') as f:
            for row in csv.reader(f):
                text_lines.append(" | ".join(row))
    except:
        with open(path, 'r', encoding='latin-1') as f:
            for row in csv.reader(f):
                text_lines.append(" | ".join(row))
    return "\n".join(text_lines)

def read_txt(path):
    with open(path, 'r', encoding='utf-8', errors='ignore') as f:
        return f.read()

def read_json(path):
    try:
        with open(path, 'r', encoding='utf-8') as f:
            data = json.load(f)
            return json.dumps(data, indent=2)
    except:
        with open(path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()

def read_xml(path):
    try:
        tree = ET.parse(path)
        root = tree.getroot()
        return ET.tostring(root, encoding='unicode')
    except:
        with open(path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()

def read_rtf(path):
    with open(path, 'r', encoding='utf-8', errors='ignore') as f:
        text = f.read()
    text = re.sub(r'\\[a-z]+\d*\s?', '', text)
    text = re.sub(r'[{}]', '', text)
    return text

def read_image(path):
    """Extract text from images using OCR"""
    try:
        img = Image.open(path)
        text = pytesseract.image_to_string(img)
        if text.strip():
            return f"[Image: {os.path.basename(path)}]\n{text}"
        else:
            return f"[Image: {os.path.basename(path)}]\n(No readable text detected in this image)"
    except Exception as e:
        return f"[Image: {os.path.basename(path)}]\n(Could not process image: {str(e)})"

# Process all files
all_chunks = []
all_ids = []
all_sources = []

for filename in os.listdir(DATA_FOLDER):
    filepath = os.path.join(DATA_FOLDER, filename)
    ext = filename.split('.')[-1].lower()
    
    print(f"Reading: {filename}")
    
    try:
        if ext == "pdf":
            text = read_pdf(filepath)
        elif ext == "pptx":
            text = read_pptx(filepath)
        elif ext == "docx":
            text = read_docx(filepath)
        elif ext in ["xlsx", "xls"]:
            text = read_xlsx(filepath)
        elif ext == "csv":
            text = read_csv(filepath)
        elif ext == "txt":
            text = read_txt(filepath)
        elif ext == "json":
            text = read_json(filepath)
        elif ext == "xml":
            text = read_xml(filepath)
        elif ext == "rtf":
            text = read_rtf(filepath)
        elif ext in ["png", "jpg", "jpeg", "bmp", "gif", "tiff"]:
            text = read_image(filepath)
        else:
            print(f"  Skipping {filename} - unsupported format")
            continue

        chunks = chunk_text(text)
        for i, chunk in enumerate(chunks):
            all_chunks.append(chunk)
            all_ids.append(f"{filename}_chunk_{i}")
            all_sources.append(filename)
            
    except Exception as e:
        print(f"  Error processing {filename}: {e}")
        continue

# Embed and store
print(f"Embedding {len(all_chunks)} chunks...")
embeddings = model.encode(all_chunks).tolist()

collection.add(
    documents=all_chunks,
    embeddings=embeddings,
    ids=all_ids,
    metadatas=[{"source": s} for s in all_sources]
)

print("✅ Done! Aquanis has learned from your documents.")
print(f"Total chunks ingested: {len(all_chunks)}")
print(f"File types: PDF, PPTX, DOCX, XLSX/XLS, CSV, TXT, JSON, XML, RTF, PNG, JPG, JPEG, BMP, GIF, TIFF")